"""
AVLpoint Daily Report — PowerPoint Renderer
Generates a branded PPTX presentation using python-pptx.
"""
from __future__ import annotations

from datetime import datetime
from pathlib import Path

from reporting.collector import ReportData
from reporting.report_config import BRAND


def _hex_to_rgb(h: str) -> tuple[int, int, int]:
    h = h.lstrip("#")
    return int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)


def render_pptx(data: ReportData, output_dir: Path) -> Path | None:
    """Generate a branded PPTX presentation. Returns path or None if python-pptx missing."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        print("[report] python-pptx not installed — skipping PPTX generation.")
        print("[report] Install with: pip install python-pptx")
        return None

    output_dir.mkdir(parents=True, exist_ok=True)
    out_path = output_dir / "report.pptx"

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    bg_color = RGBColor(*_hex_to_rgb(BRAND["bg"]))
    fg_color = RGBColor(*_hex_to_rgb(BRAND["fg"]))
    arc_color = RGBColor(*_hex_to_rgb(BRAND["arc"]))
    sec_color = RGBColor(*_hex_to_rgb(BRAND["fg_secondary"]))
    ok_color = RGBColor(*_hex_to_rgb(BRAND["ok"]))
    warn_color = RGBColor(*_hex_to_rgb(BRAND["warn"]))
    danger_color = RGBColor(*_hex_to_rgb(BRAND["danger"]))
    surface_color = RGBColor(*_hex_to_rgb(BRAND["surface"]))

    health_colors = {"healthy": ok_color, "warning": warn_color, "critical": danger_color}

    def set_bg(slide):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = bg_color

    def add_text(slide, left, top, width, height, text, font_size=14, color=fg_color, bold=False, alignment=PP_ALIGN.LEFT, font_name="Inter"):
        from pptx.util import Inches, Pt
        txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name
        p.alignment = alignment
        return txBox

    # ── Slide 1: Title ──
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    set_bg(slide1)

    add_text(slide1, 1.5, 1.5, 10, 1, "AVLpoint", 48, arc_color, True, font_name="Arial Black")
    add_text(slide1, 1.5, 2.5, 10, 0.8, "Daily Platform Report", 32, fg_color, False, font_name="Arial")

    try:
        dt = datetime.strptime(data.report_date, "%Y-%m-%d")
        date_str = dt.strftime("%B %d, %Y")
    except ValueError:
        date_str = data.report_date
    add_text(slide1, 1.5, 3.8, 10, 0.5, date_str, 20, sec_color, font_name="Consolas")
    add_text(slide1, 1.5, 4.5, 10, 0.5, f"Generated {data.generated_at}", 12, sec_color, font_name="Consolas")

    # Health indicator
    health_label = {"healthy": "🟢 All Systems Operational", "warning": "🟡 Issues Detected", "critical": "🔴 Critical — Attention Required"}
    add_text(slide1, 1.5, 5.5, 10, 0.6, health_label.get(data.overall_health.status, "Unknown"), 22, health_colors.get(data.overall_health.status, fg_color), True)

    # ── Slide 2: Executive Summary (KPIs) ──
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide2)
    add_text(slide2, 0.8, 0.4, 8, 0.6, "Executive Summary", 28, fg_color, True, font_name="Arial")

    kpis = [
        (f"{data.database.combined_total:,}", "Total Vendors", f"{data.deltas.combined_delta:+,}" if data.deltas.has_previous else ""),
        (f"{data.database.us_verified:,}", "Verified Profiles", f"{data.deltas.us_verified_delta:+,}" if data.deltas.has_previous else ""),
        (f"{data.database.us_enrichment_pct}%", "Enrichment Rate", ""),
        (f"{data.pipeline.runs_today}", "Pipeline Runs Today", ""),
        (f"{data.pipeline.total_enriched_today:,}", "Enriched Today", ""),
        (f"{data.website.avg_response_ms:.0f}ms", "Avg Response Time", ""),
    ]

    for i, (val, label, delta) in enumerate(kpis):
        col = i % 3
        row = i // 3
        x = 0.8 + col * 4.0
        y = 1.4 + row * 2.8

        # Card background
        shape = slide2.shapes.add_shape(
            1, Inches(x), Inches(y), Inches(3.5), Inches(2.2)  # MSO_SHAPE.RECTANGLE
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = surface_color
        shape.line.fill.background()

        add_text(slide2, x + 0.3, y + 0.3, 3, 0.8, val, 36, fg_color, True, font_name="Consolas")
        add_text(slide2, x + 0.3, y + 1.2, 3, 0.4, label, 12, sec_color, font_name="Arial")
        if delta:
            delta_color = ok_color if delta.startswith("+") and delta != "+0" else sec_color
            add_text(slide2, x + 0.3, y + 1.6, 3, 0.3, delta, 13, delta_color, font_name="Consolas")

    # ── Slide 3: Database Health ──
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide3)
    add_text(slide3, 0.8, 0.4, 8, 0.6, "Database Health", 28, fg_color, True, font_name="Arial")

    db_metrics = [
        ("U.S. Total Records", f"{data.database.us_total:,}"),
        ("U.S. Verified", f"{data.database.us_verified:,}"),
        ("U.S. With Description", f"{data.database.us_with_description:,}"),
        ("U.S. With Email", f"{data.database.us_with_email:,}"),
        ("U.S. With Website", f"{data.database.us_with_website:,}"),
        ("International Total", f"{data.database.intl_total:,}"),
        ("Combined Total", f"{data.database.combined_total:,}"),
        ("DB Size", f"{data.database.db_size_mb} MB"),
    ]

    for i, (label, value) in enumerate(db_metrics):
        y = 1.3 + i * 0.6
        add_text(slide3, 1.2, y, 4, 0.5, label, 14, sec_color)
        add_text(slide3, 6.0, y, 3, 0.5, value, 14, fg_color, True, PP_ALIGN.RIGHT, "Consolas")

    # ── Slide 4: Pipeline Performance ──
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide4)
    add_text(slide4, 0.8, 0.4, 8, 0.6, "Pipeline Performance", 28, fg_color, True, font_name="Arial")

    pipe_metrics = [
        ("Runs Today", str(data.pipeline.runs_today)),
        ("Vendors Enriched", f"{data.pipeline.total_enriched_today:,}"),
        ("Records Discovered", f"{data.pipeline.total_discovered_today:,}"),
        ("Errors", str(data.pipeline.total_errors_today)),
        ("Scheduler State", data.pipeline.scheduler_state or "N/A"),
        ("Last Run", data.pipeline.scheduler_last_run or "N/A"),
        ("Next Run", data.pipeline.scheduler_next_run or "N/A"),
        ("Backup Status", data.pipeline.backup_status or "N/A"),
    ]

    for i, (label, value) in enumerate(pipe_metrics):
        y = 1.3 + i * 0.6
        add_text(slide4, 1.2, y, 4, 0.5, label, 14, sec_color)
        val_color = danger_color if (label == "Errors" and int(value) > 5) else fg_color
        add_text(slide4, 6.0, y, 3, 0.5, value, 14, val_color, True, PP_ALIGN.RIGHT, "Consolas")

    # ── Slide 5: Website Status ──
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide5)
    add_text(slide5, 0.8, 0.4, 8, 0.6, "Website Status", 28, fg_color, True, font_name="Arial")

    for i, ep in enumerate(data.website.endpoints):
        y = 1.3 + i * 0.55
        status_icon = "✓" if ep.ok else "✗"
        add_text(slide5, 1.2, y, 3, 0.45, ep.name, 13, sec_color)
        add_text(slide5, 4.5, y, 1, 0.45, ep.path, 12, sec_color, font_name="Consolas")
        add_text(slide5, 7.0, y, 0.8, 0.45, status_icon, 16, ok_color if ep.ok else danger_color, True)
        add_text(slide5, 8.0, y, 1.5, 0.45, f"{ep.response_ms}ms", 13, fg_color, font_name="Consolas")

    # ── Slide 6: Issues & Recommendations ──
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide6)
    add_text(slide6, 0.8, 0.4, 8, 0.6, "Issues & Recommendations", 28, fg_color, True, font_name="Arial")

    add_text(slide6, 0.8, 1.2, 4, 0.4, "ACTIVE ISSUES", 12, sec_color, True)
    for i, issue in enumerate(data.issues[:5]):
        color = danger_color if "CRITICAL" in issue else (warn_color if "WARNING" in issue else ok_color)
        add_text(slide6, 1.2, 1.7 + i * 0.55, 5, 0.45, issue, 13, color)

    add_text(slide6, 7.0, 1.2, 4, 0.4, "RECOMMENDATIONS", 12, sec_color, True)
    for i, rec in enumerate(data.recommendations[:5]):
        add_text(slide6, 7.4, 1.7 + i * 0.55, 5, 0.45, f"→ {rec}", 12, sec_color)

    prs.save(str(out_path))
    return out_path
